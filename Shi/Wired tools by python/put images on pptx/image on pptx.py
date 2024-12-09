from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image  # For reading image dimensions
import os

# Input: List of folders containing images
folders = [
    r"D:/Chrome/2024-10-19 Hela cells D-Label/2024-10-19 Hela cells D-Label/D-Ace/1-CD subtracted",
    r"D:\Chrome\2024-10-19 Hela cells D-Label\2024-10-19 Hela cells D-Label\D-Ace\2-CD subtracted",
    r"D:\Chrome\2024-10-19 Hela cells D-Label\2024-10-19 Hela cells D-Label\D-Glutamine\1-CD subtracted",
    r"D:\Chrome\2024-10-19 Hela cells D-Label\2024-10-19 Hela cells D-Label\D-Glutamine\2-CD subtracted",
    r"D:\Chrome\2024-10-19 Hela cells D-Label\2024-10-19 Hela cells D-Label\D-Met\1-CD subtracted",
    r"D:\Chrome\2024-10-19 Hela cells D-Label\2024-10-19 Hela cells D-Label\D-Met\2-CD subtracted",
    r"D:\Chrome\2024-10-19 Hela cells D-Label\2024-11-15 Hela Cells D-Label\Choline\1-CD subtracted",
    r"D:\Chrome\2024-10-19 Hela cells D-Label\2024-11-15 Hela Cells D-Label\Leucine\1-CD subtracted",
    r"D:\Chrome\2024-10-19 Hela cells D-Label\2024-11-15 Hela Cells D-Label\L-Serine\1-CD subtracted",
    r"D:\Chrome\2024-10-19 Hela cells D-Label\2024-11-15 Hela Cells D-Label\Lysine\1-CD subtracted",
    r"D:\Chrome\2024-10-19 Hela cells D-Label\2024-11-15 Hela Cells D-Label\Tyrosine\1-CD subtracted"
]

# Output: Directory to save the PowerPoint
output_dir = "D:/Chrome/2024-10-19 Hela cells D-Label"
os.makedirs(output_dir, exist_ok=True)  # Create the directory if it doesn't exist
output_pptx = os.path.join(output_dir, "output_presentation.pptx")

# Create a new PowerPoint presentation
presentation = Presentation()

# Slide dimensions
slide_width = Inches(10)  # Standard slide width
slide_height = Inches(7.5)  # Standard slide height
margin = Inches(0.5)  # Margin around the edges

# Loop through each folder
for folder in folders:
    # Get all image files in the folder
    image_files = sorted([f for f in os.listdir(folder) if f.endswith('.png') or f.endswith('.jpg')])
    if not image_files:
        continue  # Skip empty folders

    # Add a title slide for the folder
    title_slide_layout = presentation.slide_layouts[0]  # Title slide layout
    title_slide = presentation.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    title.text = f"Images from: {folder}"  # Add folder path as the title
    subtitle = title_slide.placeholders[1]
    subtitle.text = f"Folder contains {len(image_files)} images."

    # Add a new slide for images
    slide_layout = presentation.slide_layouts[6]  # Blank slide layout
    slide = presentation.slides.add_slide(slide_layout)

    # Calculate grid dimensions
    images_per_row = 7  # Customize as needed
    available_width = slide_width - (2 * margin)
    available_height = slide_height - (2 * margin)
    rows = (len(image_files) + images_per_row - 1) // images_per_row  # Calculate total rows needed
    grid_cell_width = available_width / images_per_row
    grid_cell_height = available_height / rows

    # Ensure the cell height is proportional to prevent overlap
    if grid_cell_height < grid_cell_width:
        grid_cell_width = grid_cell_height
    else:
        grid_cell_height = grid_cell_width

    # Position and add images
    x = margin
    y = margin

    for i, img_file in enumerate(image_files):
        img_path = os.path.join(folder, img_file)

        # Get image dimensions using PIL
        with Image.open(img_path) as img:
            img_width, img_height = img.size
            img_ratio = img_width / img_height

        # Scale image while maintaining aspect ratio
        if img_ratio > 1:  # Wider than tall
            scaled_width = grid_cell_width
            scaled_height = scaled_width / img_ratio
        else:  # Taller than wide
            scaled_height = grid_cell_height
            scaled_width = scaled_height * img_ratio

        # Adjust position to center image in the grid cell
        x_offset = (grid_cell_width - scaled_width) / 2
        y_offset = (grid_cell_height - scaled_height) / 2

        # Add the image to the slide
        slide.shapes.add_picture(
            img_path,
            x + x_offset,
            y + y_offset,
            width=scaled_width,
            height=scaled_height
        )

        # Update position for next image
        x += grid_cell_width
        if (i + 1) % images_per_row == 0:  # Move to the next row
            x = margin
            y += grid_cell_height

# Save the presentation
presentation.save(output_pptx)
print(f"Presentation saved at {output_pptx}")
